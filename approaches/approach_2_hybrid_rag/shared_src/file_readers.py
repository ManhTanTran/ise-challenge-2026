"""Robust readers for files in the multimodal data lake."""

from __future__ import annotations

import contextlib
import json
import logging
import os
import re
import sqlite3
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

import pandas as pd

from .utils import (
    DEFAULT_ENCODINGS,
    detect_mime,
    dump_json,
    ensure_dir,
    load_json,
    normalize_spaces,
    read_text_with_fallback,
    safe_relative_path,
    stable_hash,
    truncate_text,
    write_text,
)

LOGGER = logging.getLogger(__name__)

TABLE_EXTENSIONS = {".csv", ".xlsx", ".xls", ".sql"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppt", ".txt", ".md", ".html", ".htm"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp"}
AUDIO_EXTENSIONS = {".m4a", ".mp3", ".wav"}
IMAGE_PARSE_SUFFIX = ".image_parse.json"


@dataclass(slots=True)
class ReadResult:
    """Common result shape returned by all readers."""

    path: str
    modality: str
    content: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    tables: Any | None = None
    error: str | None = None

    def to_dict(self, include_tables: bool = False) -> dict[str, Any]:
        data: dict[str, Any] = {
            "path": self.path,
            "modality": self.modality,
            "content": self.content,
            "metadata": self.metadata,
            "error": self.error,
        }
        if include_tables:
            data["tables"] = self.tables
        return data


def modality_for_path(path: str | Path) -> str:
    """Map an extension to a coarse modality."""

    extension = Path(path).suffix.lower()
    if extension in TABLE_EXTENSIONS:
        return "table"
    if extension in DOCUMENT_EXTENSIONS:
        return "document"
    if extension in IMAGE_EXTENSIONS:
        return "image"
    if extension in AUDIO_EXTENSIONS:
        return "audio"
    return "unknown"


def cache_path_for_file(
    path: str | Path,
    *,
    data_lake_dir: str | Path | None,
    cache_dir: str | Path,
) -> Path:
    """Return the extraction cache path for a data lake file."""

    source = Path(path)
    relative = safe_relative_path(source, data_lake_dir) if data_lake_dir else source.name
    cache_name = f"{stable_hash(relative)}_{source.stem[:80]}.txt"
    return Path(cache_dir) / cache_name


def read_file(
    path: str | Path,
    *,
    cache_dir: str | Path | None = None,
    data_lake_dir: str | Path | None = None,
    use_cache: bool = True,
) -> ReadResult:
    """Read a file without allowing a single bad file to crash the pipeline."""

    file_path = Path(path)
    modality = modality_for_path(file_path)
    if not file_path.exists():
        return ReadResult(str(file_path), modality, error="File does not exist.")

    cache_path: Path | None = None
    if cache_dir is not None:
        cache_path = cache_path_for_file(file_path, data_lake_dir=data_lake_dir, cache_dir=cache_dir)
        if use_cache and cache_path.exists():
            content = read_text_with_fallback(cache_path)
            metadata = {"cache_hit": True, "extracted_text_path": str(cache_path)}
            if modality == "image":
                parse_path = image_parse_cache_path(cache_path)
                parsed = load_json(parse_path, default={})
                if isinstance(parsed, dict):
                    metadata.update(_image_parse_metadata(parsed, parse_path))
            reader = _reader_for_extension(file_path.suffix.lower())
            if modality == "table" and reader is not None:
                with contextlib.suppress(Exception):
                    result = reader(file_path)
                    result.content = content
                    result.metadata.update(metadata)
                    return result
            return ReadResult(str(file_path), modality, content=content, metadata=metadata)

    reader = _reader_for_extension(file_path.suffix.lower())
    if reader is None:
        result = ReadResult(str(file_path), "unknown", error="Unsupported file type.")
    else:
        try:
            if modality == "image":
                result = read_image(
                    file_path,
                    parse_cache_path=image_parse_cache_path(cache_path) if cache_path else None,
                    use_cache=use_cache,
                )
            else:
                result = reader(file_path)
        except Exception as exc:  # pragma: no cover - defensive boundary
            LOGGER.exception("Failed to read %s", file_path)
            result = ReadResult(str(file_path), modality, error=str(exc))

    if cache_path is not None and result.content and result.error is None:
        write_text(cache_path, result.content)
        result.metadata["extracted_text_path"] = str(cache_path)
    return result


def read_csv(path: Path) -> ReadResult:
    """Read CSV with encoding fallback."""

    last_error: Exception | None = None
    dataframe: pd.DataFrame | None = None
    used_encoding = ""
    for encoding in DEFAULT_ENCODINGS:
        try:
            dataframe = pd.read_csv(path, encoding=encoding)
            used_encoding = encoding
            break
        except UnicodeDecodeError as exc:
            last_error = exc
        except pd.errors.ParserError:
            dataframe = pd.read_csv(path, encoding=encoding, sep=None, engine="python")
            used_encoding = encoding
            break
    if dataframe is None:
        raise RuntimeError(f"Could not read CSV: {last_error}")
    content = dataframe.head(50).to_csv(index=False)
    return ReadResult(
        str(path),
        "table",
        content=content,
        metadata={
            "columns": [str(col) for col in dataframe.columns],
            "shape": list(dataframe.shape),
            "encoding": used_encoding,
        },
    )


def read_excel(path: Path) -> ReadResult:
    """Read Excel sheet names, columns, and a small preview without loading full workbooks."""

    if path.suffix.lower() == ".xlsx":
        previews: list[str] = []
        sheet_meta: dict[str, Any] = {}
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(path, read_only=True, data_only=True)
            for worksheet in workbook.worksheets:
                rows = [
                    ["" if value is None else value for value in row]
                    for row in worksheet.iter_rows(max_row=31, values_only=True)
                ]
                header = [str(value) for value in rows[0]] if rows else []
                preview_frame = pd.DataFrame(rows[1:], columns=header) if header else pd.DataFrame()
                sheet_meta[worksheet.title] = {
                    "columns": header,
                    "shape": [worksheet.max_row or 0, worksheet.max_column or len(header)],
                }
                previews.append(f"Sheet: {worksheet.title}\n{preview_frame.to_csv(index=False)}")
            workbook.close()
            return ReadResult(
                str(path),
                "table",
                content="\n\n".join(previews),
                metadata={
                    "sheet_names": list(sheet_meta.keys()),
                    "sheets": sheet_meta,
                    "preview_rows_per_sheet": 30,
                },
            )
        except Exception:
            LOGGER.exception("Fast Excel preview failed for %s; falling back to pandas", path)

    sheets = pd.read_excel(path, sheet_name=None, nrows=30)
    previews = []
    sheet_meta = {}
    for sheet_name, dataframe in sheets.items():
        sheet_meta[sheet_name] = {
            "columns": [str(col) for col in dataframe.columns],
            "shape": [len(dataframe), len(dataframe.columns)],
        }
        previews.append(f"Sheet: {sheet_name}\n{dataframe.to_csv(index=False)}")
    return ReadResult(
        str(path),
        "table",
        content="\n\n".join(previews),
        metadata={
            "sheet_names": list(sheets.keys()),
            "sheets": sheet_meta,
        },
    )


def read_sql(path: Path) -> ReadResult:
    """Read SQL text and try to infer sqlite table names."""

    text = read_text_with_fallback(path)
    table_names: list[str] = []
    with contextlib.suppress(Exception):
        with sqlite3.connect(":memory:") as connection:
            connection.executescript(text)
            rows = connection.execute(
                "SELECT name FROM sqlite_master WHERE type='table' ORDER BY name"
            ).fetchall()
            table_names = [row[0] for row in rows]
    return ReadResult(
        str(path),
        "table",
        content=text,
        metadata={"table_names": table_names},
    )


def read_text(path: Path) -> ReadResult:
    """Read plain text and markdown files."""

    return ReadResult(str(path), "document", content=read_text_with_fallback(path))


def read_html(path: Path) -> ReadResult:
    """Extract visible text from HTML."""

    from bs4 import BeautifulSoup

    html = read_text_with_fallback(path)
    soup = BeautifulSoup(html, _html_parser())
    for node in soup(["script", "style", "noscript"]):
        node.decompose()
    text = soup.get_text(separator=" ")
    return ReadResult(str(path), "document", content=truncate_text(text, 200000))


def read_pdf(path: Path) -> ReadResult:
    """Extract text from PDF with PyMuPDF, falling back to pdfplumber."""

    try:
        import fitz

        pieces: list[str] = []
        with fitz.open(path) as document:
            for page_number, page in enumerate(document, start=1):
                pieces.append(f"\n[Page {page_number}]\n{page.get_text()}")
        return ReadResult(str(path), "document", content="\n".join(pieces))
    except Exception as first_error:
        LOGGER.debug("PyMuPDF failed for %s: %s", path, first_error)
        import pdfplumber

        pieces = []
        with pdfplumber.open(path) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                pieces.append(f"\n[Page {page_number}]\n{page.extract_text() or ''}")
        return ReadResult(str(path), "document", content="\n".join(pieces))


def read_docx(path: Path) -> ReadResult:
    """Extract paragraphs and table text from DOCX."""

    from docx import Document

    document = Document(path)
    pieces = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            pieces.append(" | ".join(cell.text.strip() for cell in row.cells))
    return ReadResult(str(path), "document", content="\n".join(pieces))


def read_pptx(path: Path) -> ReadResult:
    """Extract text from PPTX slides."""

    from pptx import Presentation

    presentation = Presentation(path)
    pieces: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            pieces.append(f"\n[Slide {index}]\n" + "\n".join(slide_text))
    return ReadResult(
        str(path),
        "document",
        content="\n".join(pieces),
        metadata={"slide_count": len(presentation.slides)},
    )


def read_ppt(path: Path) -> ReadResult:
    """Best-effort legacy PPT extraction via LibreOffice conversion."""

    executable = _find_executable("soffice") or _find_executable("libreoffice")
    if executable is None:
        return ReadResult(
            str(path),
            "document",
            content="",
            metadata={"needs_conversion": True},
            error="Legacy .ppt requires LibreOffice/soffice conversion to .pptx.",
        )
    with tempfile.TemporaryDirectory() as temp_dir:
        subprocess.run(
            [
                executable,
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                temp_dir,
                str(path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        converted = Path(temp_dir) / f"{path.stem}.pptx"
        if not converted.exists():
            return ReadResult(str(path), "document", error="PPT conversion produced no file.")
        result = read_pptx(converted)
        result.path = str(path)
        result.metadata["converted_from_ppt"] = True
        return result


def image_parse_cache_path(text_cache_path: str | Path) -> Path:
    """Return the structured image parse cache path for an extracted text path."""

    return Path(f"{Path(text_cache_path)}{IMAGE_PARSE_SUFFIX}")


def read_image(
    path: Path,
    *,
    parse_cache_path: str | Path | None = None,
    use_cache: bool = True,
) -> ReadResult:
    """Parse image text into plain text plus structured OCR metadata."""

    cache_path = Path(parse_cache_path) if parse_cache_path else None
    if use_cache and cache_path and cache_path.exists():
        parsed = load_json(cache_path, default={})
        if isinstance(parsed, dict):
            return ReadResult(
                str(path),
                "image",
                content=str(parsed.get("plain_text", "")),
                metadata=_image_parse_metadata(parsed, cache_path),
            )

    try:
        parsed = _local_image_parse(path)
        if _should_use_vlm_image_parse(parsed):
            vlm_parse = _vlm_image_parse(path)
            if vlm_parse:
                parsed = _merge_image_parses(parsed, vlm_parse)
        if cache_path:
            dump_json(parsed, cache_path)
        return ReadResult(
            str(path),
            "image",
            content=str(parsed.get("plain_text", "")),
            metadata=_image_parse_metadata(parsed, cache_path),
        )
    except Exception as exc:
        return ReadResult(
            str(path),
            "image",
            content="",
            metadata={"needs_vision_model": True},
            error=f"OCR unavailable or failed: {exc}",
        )


def _local_image_parse(path: Path) -> dict[str, Any]:
    """Run local Tesseract OCR with several preprocessing/layout attempts."""

    from PIL import Image
    import pytesseract

    tesseract = _find_executable("tesseract")
    if tesseract:
        pytesseract.pytesseract.tesseract_cmd = tesseract

    image = Image.open(path)
    variants = _ocr_image_variants(image)
    language = _tesseract_language(pytesseract)
    timeout = float(os.getenv("ISE_OCR_TIMEOUT_SECONDS", "4"))
    psm_modes = [
        mode.strip()
        for mode in os.getenv("ISE_OCR_PSM_MODES", "6,11").split(",")
        if mode.strip()
    ]
    configs = [f"--oem 3 --psm {mode}" for mode in psm_modes]

    attempts: list[dict[str, Any]] = []
    last_error = ""
    for variant_name, variant in variants:
        for config in configs:
            try:
                data = pytesseract.image_to_data(
                    variant,
                    lang=language,
                    config=config,
                    output_type=pytesseract.Output.DICT,
                    timeout=timeout,
                )
                parsed = _parse_tesseract_data(data)
                parsed.update(
                    {
                        "engine": "pytesseract",
                        "language": language,
                        "ocr_config": config,
                        "preprocess": variant_name,
                        "size": list(image.size),
                    }
                )
                attempts.append(parsed)
            except Exception as exc:
                last_error = str(exc)

    if not attempts:
        raise RuntimeError(last_error or "Tesseract produced no OCR attempts.")

    best = max(attempts, key=_score_image_parse)
    best["attempt_count"] = len(attempts)
    best["needs_vision_model"] = _should_use_vlm_image_parse(best, ignore_env=True)
    return best


def _ocr_image_variants(image: Any) -> list[tuple[str, Any]]:
    from PIL import ImageEnhance, ImageFilter, ImageOps

    rgb = image.convert("RGB")
    max_side = max(rgb.size)
    if max_side < 1600:
        scale = min(3.0, 1600 / max_side)
        rgb = rgb.resize((int(rgb.width * scale), int(rgb.height * scale)))
    elif max_side > 2600:
        scale = 2600 / max_side
        rgb = rgb.resize((int(rgb.width * scale), int(rgb.height * scale)))

    gray = ImageOps.grayscale(rgb)
    enhanced = ImageEnhance.Contrast(gray).enhance(1.8)
    enhanced = ImageEnhance.Sharpness(enhanced).enhance(1.4)
    denoised = enhanced.filter(ImageFilter.MedianFilter(size=3))
    thresholded = denoised.point(lambda pixel: 255 if pixel > 180 else 0)
    return [
        ("gray_enhanced", denoised),
        ("threshold", thresholded),
    ]


def _tesseract_language(pytesseract_module: Any) -> str:
    preferred = [part.strip() for part in os.getenv("ISE_OCR_LANG", "vie+eng").split("+") if part.strip()]
    try:
        available = set(pytesseract_module.get_languages(config=""))
    except Exception:
        available = set()
    selected = [language for language in preferred if not available or language in available]
    if selected:
        return "+".join(selected)
    return "eng"


def _parse_tesseract_data(data: dict[str, list[Any]]) -> dict[str, Any]:
    words: list[dict[str, Any]] = []
    total_confidence = 0.0
    confidence_count = 0
    row_count = len(data.get("text", []))

    for index in range(row_count):
        text = _repair_mojibake(normalize_spaces(data.get("text", [""])[index]))
        confidence = _safe_confidence(data.get("conf", ["-1"])[index])
        if not text or confidence < 0:
            continue
        word = {
            "text": text,
            "confidence": confidence,
            "left": _safe_int(data.get("left", [0])[index]),
            "top": _safe_int(data.get("top", [0])[index]),
            "width": _safe_int(data.get("width", [0])[index]),
            "height": _safe_int(data.get("height", [0])[index]),
            "block_num": _safe_int(data.get("block_num", [0])[index]),
            "par_num": _safe_int(data.get("par_num", [0])[index]),
            "line_num": _safe_int(data.get("line_num", [0])[index]),
            "word_num": _safe_int(data.get("word_num", [0])[index]),
        }
        words.append(word)
        total_confidence += confidence
        confidence_count += 1

    lines = _ocr_lines(words)
    for line in lines:
        line["text"] = _repair_mojibake(line.get("text", ""))
    blocks = _ocr_blocks(lines)
    plain_text = _repair_mojibake("\n".join(line["text"] for line in lines if line.get("text")))
    return {
        "plain_text": plain_text,
        "blocks": blocks,
        "lines": lines,
        "words": words,
        "tables": _table_candidates_from_lines(lines),
        "key_values": {},
        "avg_confidence": round(total_confidence / confidence_count, 2) if confidence_count else 0.0,
    }


def _ocr_lines(words: list[dict[str, Any]]) -> list[dict[str, Any]]:
    grouped: dict[tuple[int, int, int], list[dict[str, Any]]] = {}
    for word in words:
        key = (
            int(word.get("block_num", 0) or 0),
            int(word.get("par_num", 0) or 0),
            int(word.get("line_num", 0) or 0),
        )
        grouped.setdefault(key, []).append(word)

    lines: list[dict[str, Any]] = []
    for key, line_words in grouped.items():
        ordered = sorted(line_words, key=lambda item: int(item.get("left", 0) or 0))
        confidences = [float(item.get("confidence", 0.0) or 0.0) for item in ordered]
        left = min(int(item.get("left", 0) or 0) for item in ordered)
        top = min(int(item.get("top", 0) or 0) for item in ordered)
        right = max(int(item.get("left", 0) or 0) + int(item.get("width", 0) or 0) for item in ordered)
        bottom = max(int(item.get("top", 0) or 0) + int(item.get("height", 0) or 0) for item in ordered)
        lines.append(
            {
                "text": normalize_spaces(" ".join(str(item.get("text", "")) for item in ordered)),
                "confidence": round(sum(confidences) / len(confidences), 2) if confidences else 0.0,
                "bbox": [left, top, right - left, bottom - top],
                "block_num": key[0],
                "par_num": key[1],
                "line_num": key[2],
            }
        )
    return sorted(lines, key=lambda item: (item["bbox"][1], item["bbox"][0]))


def _ocr_blocks(lines: list[dict[str, Any]]) -> list[dict[str, Any]]:
    grouped: dict[int, list[dict[str, Any]]] = {}
    for line in lines:
        grouped.setdefault(int(line.get("block_num", 0) or 0), []).append(line)
    blocks: list[dict[str, Any]] = []
    for block_num, block_lines in grouped.items():
        text = "\n".join(line.get("text", "") for line in block_lines if line.get("text"))
        if not text.strip():
            continue
        confidences = [float(line.get("confidence", 0.0) or 0.0) for line in block_lines]
        blocks.append(
            {
                "block_num": block_num,
                "text": text,
                "confidence": round(sum(confidences) / len(confidences), 2) if confidences else 0.0,
                "line_count": len(block_lines),
            }
        )
    return sorted(blocks, key=lambda item: item["block_num"])


def _table_candidates_from_lines(lines: list[dict[str, Any]]) -> list[dict[str, Any]]:
    row_texts = [line.get("text", "") for line in lines if len(str(line.get("text", "")).split()) >= 3]
    if len(row_texts) < 3:
        return []
    rows = [[cell for cell in re.split(r"\s{2,}|\s+\|\s+", row) if cell] or [row] for row in row_texts]
    return [{"rows": rows, "row_count": len(rows), "source": "ocr_lines"}]


def _score_image_parse(parsed: dict[str, Any]) -> float:
    text = str(parsed.get("plain_text", ""))
    return float(parsed.get("avg_confidence", 0.0) or 0.0) + min(len(text) / 100.0, 25.0)


def _should_use_vlm_image_parse(parsed: dict[str, Any], *, ignore_env: bool = False) -> bool:
    if not ignore_env and os.getenv("ISE_IMAGE_PARSE_VLM", "auto").lower() in {"0", "false", "off", "never"}:
        return False
    text = str(parsed.get("plain_text", ""))
    threshold = float(os.getenv("ISE_OCR_MIN_CONFIDENCE", "55"))
    word_count = len(parsed.get("words", []) if isinstance(parsed.get("words", []), list) else [])
    line_count = len(parsed.get("lines", []) if isinstance(parsed.get("lines", []), list) else [])
    looks_like_document = word_count >= 5 or line_count >= 3 or len(normalize_spaces(text)) >= 80
    low_confidence = float(parsed.get("avg_confidence", 0.0) or 0.0) < threshold
    short_document_text = len(normalize_spaces(text)) < 25 and looks_like_document
    return looks_like_mojibake(text) or (looks_like_document and low_confidence) or short_document_text


def looks_like_mojibake(text: str) -> bool:
    """Detect common UTF-8-as-Latin mojibake that appears in weak OCR output."""

    return _mojibake_score(text) >= 2


def _repair_mojibake(text: str) -> str:
    if not text or not looks_like_mojibake(text):
        return text
    original_score = _mojibake_score(text)
    candidates = []
    for encoding in ("cp1252", "latin-1"):
        with contextlib.suppress(Exception):
            if encoding == "cp1252":
                raw = _bytes_from_cp1252_mojibake(text)
            else:
                raw = text.encode("latin-1")
            candidates.append(raw.decode("utf-8"))
    if not candidates:
        return text
    best = min(candidates, key=_mojibake_score)
    return best if _mojibake_score(best) < original_score else text


def _bytes_from_cp1252_mojibake(text: str) -> bytes:
    raw = bytearray()
    for char in text:
        codepoint = ord(char)
        if codepoint <= 255:
            raw.append(codepoint)
        else:
            raw.extend(char.encode("cp1252"))
    return bytes(raw)


def _mojibake_score(text: str) -> int:
    if not text:
        return 0
    suspicious = len(re.findall(r"(Ã.|Â.|Ä.|Å.|Æ.|áº|á»|â€|â€™|â€œ|â€)", text))
    controls = sum(1 for char in text if 0x80 <= ord(char) <= 0x9F)
    replacements = text.count("�")
    return suspicious * 2 + controls + replacements * 3


def _vlm_image_parse(path: Path) -> dict[str, Any]:
    try:
        from .config import VISION_CACHE_DIR
        from .llm_client import answer_image_from_file, has_llm

        cache_path = _vlm_image_parse_cache_path(path, VISION_CACHE_DIR)
        cached = load_json(cache_path, default=None)
        if isinstance(cached, dict):
            return cached
        if not has_llm():
            return {}
        prompt = """
Parse this image for a retrieval and question-answering pipeline.
Return only JSON with keys:
- plain_text: clean text visible in the image, preserving Vietnamese accents when present
- tables: array of tables, each with rows as arrays of cell strings
- key_values: object of important labels/values
- confidence: number from 0 to 1
Do not answer any hidden question. Only transcribe or structure visible image content.
""".strip()
        raw = answer_image_from_file(prompt, path)
        parsed = _parse_jsonish_object(raw)
        if not parsed:
            parsed = {"plain_text": raw}
        result = {
            "plain_text": normalize_spaces(parsed.get("plain_text", "")),
            "tables": parsed.get("tables", []) if isinstance(parsed.get("tables", []), list) else [],
            "key_values": parsed.get("key_values", {}) if isinstance(parsed.get("key_values", {}), dict) else {},
            "confidence": float(parsed.get("confidence", 0.0) or 0.0),
            "engine": "vlm",
        }
        dump_json(result, cache_path)
        return result
    except Exception as exc:
        LOGGER.warning("VLM image parse failed for %s: %s", path, exc)
        return {}


def _vlm_image_parse_cache_path(path: Path, cache_dir: str | Path) -> Path:
    stat_key = ""
    if path.exists():
        stat = path.stat()
        stat_key = f":{stat.st_size}:{int(stat.st_mtime)}"
    key = stable_hash(f"image_parse:{path.as_posix()}:{stat_key}")
    return ensure_dir(cache_dir) / f"{key}_image_parse.json"


def _merge_image_parses(local_parse: dict[str, Any], vlm_parse: dict[str, Any]) -> dict[str, Any]:
    local_text = normalize_spaces(local_parse.get("plain_text", ""))
    vlm_text = normalize_spaces(vlm_parse.get("plain_text", ""))
    pieces = [piece for piece in [local_text, vlm_text] if piece]
    merged_text = "\n\n".join(dict.fromkeys(pieces))
    merged = dict(local_parse)
    merged["plain_text"] = merged_text
    merged["tables"] = local_parse.get("tables", []) or vlm_parse.get("tables", [])
    merged["key_values"] = vlm_parse.get("key_values", {}) or local_parse.get("key_values", {})
    merged["vlm_parse"] = vlm_parse
    merged["engine"] = "pytesseract+vlm"
    merged["needs_vision_model"] = False
    return merged


def _image_parse_metadata(parsed: dict[str, Any], parse_path: str | Path | None = None) -> dict[str, Any]:
    metadata = {
        "ocr_engine": parsed.get("engine", "pytesseract"),
        "ocr_confidence": parsed.get("avg_confidence", parsed.get("confidence", 0.0)),
        "ocr_language": parsed.get("language", ""),
        "needs_vision_model": parsed.get("needs_vision_model", False),
    }
    if parse_path:
        metadata["image_parse_path"] = str(parse_path)
    return metadata


def _parse_jsonish_object(text: str) -> dict[str, Any]:
    match = re.search(r"\{[\s\S]*\}", text or "")
    if not match:
        return {}
    try:
        parsed = json.loads(match.group(0))
    except json.JSONDecodeError:
        return {}
    return parsed if isinstance(parsed, dict) else {}


def _safe_confidence(value: Any) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return -1.0


def _safe_int(value: Any) -> int:
    try:
        return int(float(value))
    except (TypeError, ValueError):
        return 0


def read_audio(path: Path) -> ReadResult:
    """Transcribe audio with Whisper if installed."""

    try:
        _ensure_tool_on_path("ffmpeg")
        import whisper

        whisper_cache = Path(
            os.getenv("WHISPER_CACHE_DIR", str(Path.cwd() / "outputs" / "whisper_cache"))
        )
        ensure_dir(whisper_cache)
        model = whisper.load_model("base", download_root=str(whisper_cache))
        transcription = model.transcribe(str(path))
        return ReadResult(
            str(path),
            "audio",
            content=str(transcription.get("text", "")).strip(),
            metadata={"transcription_engine": "whisper-base"},
        )
    except Exception as exc:
        return ReadResult(
            str(path),
            "audio",
            content="",
            metadata={"needs_audio_transcription": True},
            error=f"Audio transcription unavailable or failed: {exc}",
        )


def load_table_file(path: str | Path) -> dict[str, pd.DataFrame]:
    """Load CSV, Excel, or SQL content into pandas dataframes."""

    file_path = Path(path)
    extension = file_path.suffix.lower()
    if extension == ".csv":
        for encoding in DEFAULT_ENCODINGS:
            try:
                return {file_path.stem: pd.read_csv(file_path, encoding=encoding)}
            except UnicodeDecodeError:
                continue
        return {file_path.stem: pd.read_csv(file_path, encoding="latin-1", encoding_errors="replace")}
    if extension in {".xlsx", ".xls"}:
        return pd.read_excel(file_path, sheet_name=None)
    if extension == ".sql":
        with sqlite3.connect(":memory:") as connection:
            connection.executescript(read_text_with_fallback(file_path))
            names = [
                row[0]
                for row in connection.execute(
                    "SELECT name FROM sqlite_master WHERE type='table' ORDER BY name"
                ).fetchall()
            ]
            return {name: pd.read_sql_query(f'SELECT * FROM "{name}"', connection) for name in names}
    raise ValueError(f"Unsupported table file: {file_path}")


def load_sqlite_connection(path: str | Path) -> sqlite3.Connection:
    """Load a SQL dump into an in-memory sqlite connection."""

    connection = sqlite3.connect(":memory:")
    connection.executescript(read_text_with_fallback(path))
    return connection


def _reader_for_extension(extension: str) -> Callable[[Path], ReadResult] | None:
    readers: dict[str, Callable[[Path], ReadResult]] = {
        ".csv": read_csv,
        ".xlsx": read_excel,
        ".xls": read_excel,
        ".sql": read_sql,
        ".txt": read_text,
        ".md": read_text,
        ".html": read_html,
        ".htm": read_html,
        ".pdf": read_pdf,
        ".docx": read_docx,
        ".pptx": read_pptx,
        ".ppt": read_ppt,
        ".jpg": read_image,
        ".jpeg": read_image,
        ".png": read_image,
        ".webp": read_image,
        ".m4a": read_audio,
        ".mp3": read_audio,
        ".wav": read_audio,
    }
    return readers.get(extension)


def _has_executable(name: str) -> bool:
    return _find_executable(name) is not None


def _find_executable(name: str) -> str | None:
    from shutil import which

    executable = which(name)
    if executable:
        return executable

    candidates = _executable_candidates(name)
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return None


def _ensure_tool_on_path(name: str) -> None:
    executable = _find_executable(name)
    if executable:
        directory = str(Path(executable).parent)
        paths = os.environ.get("PATH", "").split(os.pathsep)
        if directory not in paths:
            os.environ["PATH"] = directory + os.pathsep + os.environ.get("PATH", "")


def _executable_candidates(name: str) -> list[Path]:
    suffix = ".exe" if os.name == "nt" and not name.lower().endswith(".exe") else ""
    executable = f"{name}{suffix}"
    candidates = [
        Path("D:/iSE challenge/tools/Tesseract-OCR") / executable,
        Path("D:/iSE challenge/tools/LibreOffice/program") / executable,
        Path("D:/iSE challenge/tools/ffmpeg/bin") / executable,
        Path("C:/Program Files/Tesseract-OCR") / executable,
        Path("C:/Program Files/LibreOffice/program") / executable,
        Path("C:/Program Files (x86)/LibreOffice/program") / executable,
    ]
    local_app_data = os.getenv("LOCALAPPDATA")
    if local_app_data:
        winget_root = Path(local_app_data) / "Microsoft" / "WinGet" / "Packages"
        if winget_root.exists():
            candidates.extend(winget_root.glob(f"**/{executable}"))
    return candidates


def _html_parser() -> str:
    """Prefer lxml, but keep HTML extraction working without it."""

    try:
        import lxml  # noqa: F401

        return "lxml"
    except ImportError:
        return "html.parser"


def file_metadata(path: str | Path, data_lake_dir: str | Path) -> dict[str, Any]:
    """Base metadata shared by the indexer."""

    file_path = Path(path)
    stat = file_path.stat()
    return {
        "filename": file_path.name,
        "relative_path": safe_relative_path(file_path, data_lake_dir),
        "absolute_path": str(file_path.resolve()),
        "extension": file_path.suffix.lower(),
        "modality": modality_for_path(file_path),
        "mime_type": detect_mime(file_path),
        "size_bytes": stat.st_size,
    }
